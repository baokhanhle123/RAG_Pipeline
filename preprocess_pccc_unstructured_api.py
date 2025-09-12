#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Preprocess PCCC legal docs (DOC/DOCX/PDF via Unstructured API; PPT/PPTX offline)
- Append to JSONL (skip duplicates by SHA1)
- Split large PDFs by page ranges
- Preserve absolute page/slide numbers
- NEW: For PPT/PPTX, parse locally using python-pptx (no Unstructured API)
- NEW: Uniform top-level identifiers for Parent-Doc retriever:
       source_sha1, source_name, source_path, source_ext, doc_id

Env:
  UNSTRUCTURED_API_URL   (default: https://api.unstructuredapp.io/general/v0/general)
  UNSTRUCTURED_API_KEY   (required for DOC/DOCX/PDF only)
"""

from __future__ import annotations
import argparse
import hashlib
import io
import json
import mimetypes
import os
from pathlib import Path
import sys
import time
from typing import Dict, Iterable, List, Optional, Tuple

import requests
from requests.adapters import HTTPAdapter
from urllib3.util.retry import Retry

# Local PPTX
try:
    from pptx import Presentation  # pip install python-pptx
except Exception:
    Presentation = None

# Optional streaming multipart
try:
    from requests_toolbelt.multipart.encoder import MultipartEncoder
except Exception:
    MultipartEncoder = None

# Optional PDF split
try:
    from pypdf import PdfReader, PdfWriter
except Exception:
    PdfReader = None
    PdfWriter = None

# ----------------- Config -----------------
SUPPORTED_EXTS = {
    ".doc", ".docx",
    ".pdf",
    ".ppt", ".pptx", ".pptm", ".pot", ".potx",
}
PPT_EXTS = {".ppt", ".pptx", ".pptm", ".pot", ".potx"}

DEFAULT_API_URL = os.environ.get("UNSTRUCTURED_API_URL", "https://api.unstructuredapp.io/general/v0/general")
API_KEY = os.environ.get("UNSTRUCTURED_API_KEY", "").strip()

TIMEOUT_CONNECT = float(os.environ.get("UNSTRUCTURED_TIMEOUT_CONNECT_S", "15"))
TIMEOUT_READ = float(os.environ.get("UNSTRUCTURED_TIMEOUT_READ_S", "120"))

MIN_UPLOAD_KBPS = float(os.environ.get("UNSTRUCTURED_MIN_UPLOAD_KBPS", "256"))
MAX_READ_TIMEOUT_S = float(os.environ.get("UNSTRUCTURED_MAX_READ_TIMEOUT_S", "1800"))
UPLOAD_SAFETY_PAD_S = float(os.environ.get("UNSTRUCTURED_UPLOAD_TIMEOUT_SAFETY", "45"))

DEFAULT_MAX_PDF_PAGES_PER_CHUNK = 60

PDF_PARAMS = {
    "strategy": "hi_res",
    "coordinates": "true",
    "include_page_breaks": "true",
    "pdf_infer_table_structure": "true",
    "languages": "vie",
}
GENERIC_PARAMS = {
    "strategy": "auto",
    "include_page_breaks": "true",
    "languages": "vie",
}

# ----------------- Helpers -----------------
def sha1_of_file(path: Path, chunk: int = 1024 * 1024) -> str:
    h = hashlib.sha1()
    with path.open("rb") as f:
        while True:
            b = f.read(chunk)
            if not b:
                break
            h.update(b)
    return h.hexdigest()

def load_existing_hashes(jsonl_path: Path) -> set:
    seen = set()
    if not jsonl_path.exists():
        return seen
    with jsonl_path.open("r", encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            if not line:
                continue
            try:
                obj = json.loads(line)
            except Exception:
                continue
            # accept both keys
            top = obj
            meta = obj.get("metadata") or {}
            fs = top.get("source_sha1") or meta.get("file_sha1")
            if fs:
                seen.add(fs)
    return seen

def detect_mime(path: Path) -> str:
    mime, _ = mimetypes.guess_type(str(path))
    return mime or "application/octet-stream"

def build_partition_params_for(path: Path) -> Dict[str, str]:
    return dict(PDF_PARAMS if path.suffix.lower()==".pdf" else GENERIC_PARAMS)

def chunk_pdf_bytes(source_pdf: Path, max_pages: int) -> Iterable[Tuple[bytes, int, int]]:
    if PdfReader is None:
        yield (source_pdf.read_bytes(), 1, -1)
        return
    reader = PdfReader(str(source_pdf))
    total = len(reader.pages)
    if total <= max_pages:
        yield (source_pdf.read_bytes(), 1, total)
        return
    start = 1
    while start <= total:
        end = min(start + max_pages - 1, total)
        writer = PdfWriter()
        for p in range(start - 1, end):
            writer.add_page(reader.pages[p])
        buf = io.BytesIO()
        writer.write(buf)
        yield (buf.getvalue(), start, end)
        start = end + 1

def _compute_upload_timeouts(file_size_bytes: int) -> Tuple[float, float]:
    connect_t = max(3.05, TIMEOUT_CONNECT)
    base_read = TIMEOUT_READ
    if file_size_bytes <= 0 or MIN_UPLOAD_KBPS <= 0:
        return (connect_t, max(base_read, 300.0))
    est_seconds = (file_size_bytes / 1024.0) / MIN_UPLOAD_KBPS
    read_t = max(base_read, est_seconds * 3.0 + UPLOAD_SAFETY_PAD_S)
    return (connect_t, min(read_t, MAX_READ_TIMEOUT_S))

def _make_session() -> requests.Session:
    s = requests.Session()
    retries = Retry(total=0, connect=0, read=0, backoff_factor=0, status_forcelist=[], allowed_methods=frozenset(["POST"]), raise_on_status=False)
    adapter = HTTPAdapter(max_retries=retries, pool_connections=4, pool_maxsize=4)
    s.mount("https://", adapter)
    s.mount("http://", adapter)
    return s

# ----------- Unstructured API posting (DOC/DOCX/PDF only) -----------
def _post_to_unstructured_bytes(file_bytes: bytes, filename: str, extra_params: Dict[str, str],
                                timeout_pair: Tuple[float, float], max_retries: int = 4) -> List[dict]:
    if not API_KEY:
        raise RuntimeError("UNSTRUCTURED_API_KEY is not set.")
    url = DEFAULT_API_URL
    headers = {"accept": "application/json", "unstructured-api-key": API_KEY}
    last_err = None
    sess = _make_session()
    for attempt in range(1, max_retries + 1):
        try:
            files = {"files": (filename, io.BytesIO(file_bytes), detect_mime(Path(filename)))}
            resp = sess.post(url, headers=headers, files=files, data=extra_params, timeout=timeout_pair)
            if resp.status_code in (200, 201):
                return resp.json()
            if resp.status_code >= 500 or resp.status_code in (408, 429):
                last_err = RuntimeError(f"HTTP {resp.status_code}: {resp.text[:300]}")
                time.sleep(min(2 ** attempt, 10))
                continue
            resp.raise_for_status()
        except Exception as e:
            last_err = e
            time.sleep(min(2 ** attempt, 10))
    raise RuntimeError(f"Failed to call Unstructured API after {max_retries} attempts: {last_err}")

def _post_to_unstructured_stream(path: Path, extra_params: Dict[str, str],
                                 timeout_pair: Tuple[float, float], max_retries: int = 4) -> List[dict]:
    if not API_KEY:
        raise RuntimeError("UNSTRUCTURED_API_KEY is not set.")
    url = DEFAULT_API_URL
    last_err = None
    sess = _make_session()
    headers = {"accept": "application/json", "unstructured-api-key": API_KEY}
    for attempt in range(1, max_retries + 1):
        try:
            if MultipartEncoder is not None:
                f = path.open("rb")
                try:
                    fields = dict(extra_params)
                    fields["files"] = (path.name, f, detect_mime(path))
                    enc = MultipartEncoder(fields=fields)
                    headers["Content-Type"] = enc.content_type
                    resp = sess.post(url, headers=headers, data=enc, timeout=timeout_pair)
                finally:
                    try: f.close()
                    except Exception: pass
            else:
                with path.open("rb") as f:
                    files = {"files": (path.name, f, detect_mime(path))}
                    resp = sess.post(url, headers=headers, files=files, data=extra_params, timeout=timeout_pair)
            if resp.status_code in (200, 201):
                return resp.json()
            if resp.status_code >= 500 or resp.status_code in (408, 429):
                last_err = RuntimeError(f"HTTP {resp.status_code}: {resp.text[:300]}")
                time.sleep(min(2 ** attempt, 10))
                continue
            resp.raise_for_status()
        except Exception as e:
            last_err = e
            time.sleep(min(2 ** attempt, 10))
            continue
    raise RuntimeError(f"Failed to call Unstructured API after {max_retries} attempts: {last_err}")

# ----------- Normalization -----------
def _add_source_top_level(obj: dict, *, file_sha1: str, file_path: Path) -> dict:
    """Ensure unified top-level identifiers for downstream Parent-Doc logic."""
    obj["source_sha1"] = file_sha1
    obj["source_name"] = file_path.name
    obj["source_path"] = str(file_path)
    obj["source_ext"]  = file_path.suffix.lower()
    # doc_id: use stem; can be replaced later after title-extraction
    obj["doc_id"] = file_path.stem
    return obj

def _normalize_element(e: dict, *, file_sha1: str, file_path: Path, abs_page_offset: int = 0) -> dict:
    meta = e.get("metadata") or {}
    page_no = meta.get("page_number")
    abs_page = None
    if isinstance(page_no, int):
        abs_page = page_no + max(abs_page_offset, 0)
    languages = meta.get("languages") or meta.get("language") or None

    normalized = {
        "element_id": e.get("element_id") or e.get("id"),
        "type": e.get("type"),
        "text": e.get("text"),
        "metadata": {
            "file_name": file_path.name,
            "file_path": str(file_path),
            "file_sha1": file_sha1,
            "source": "unstructured_api",
            "mime_type": detect_mime(file_path),
            "page_number": page_no,
            "abs_page_number": abs_page,
            "slide_number": page_no if file_path.suffix.lower() in PPT_EXTS else None,
            "languages": languages,
            "coordinates": meta.get("coordinates"),
            "text_as_html": meta.get("text_as_html"),
            "parent_id": meta.get("parent_id"),
            "section": meta.get("section"),
            "category_depth": meta.get("category_depth"),
        },
    }
    normalized["metadata"] = {k: v for k, v in normalized["metadata"].items() if v is not None}
    return _add_source_top_level(normalized, file_sha1=file_sha1, file_path=file_path)

# ----------- PPTX offline parsing -----------
def _pptx_extract_markdown_table(tbl) -> str:
    # Very simple conversion: first row -> header; others -> body
    rows = []
    for r in tbl.rows:
        row = []
        for c in r.cells:
            txt = (c.text or "").strip().replace("\n", " ")
            rows.append(None) if False else None  # placeholder
            row.append(txt)
        rows.append("| " + " | ".join([x or "" for x in row]) + " |")
    if not rows:
        return ""
    # build header separator from header length
    first = rows[0]
    ncol = first.count("|") - 1 if first else 0
    if ncol <= 0:
        return "\n".join(rows)
    sep = "|" + "|".join([" --- "]*ncol) + "|"
    return "\n".join([rows[0], sep] + rows[1:])

def _pptx_process(path: Path, file_sha1: str) -> List[dict]:
    if Presentation is None:
        raise RuntimeError("python-pptx is not installed, cannot process PPT/PPTX offline.")
    prs = Presentation(str(path))
    out: List[dict] = []
    for s_idx, slide in enumerate(prs.slides, start=1):
        # Title (if any)
        title_txt = ""
        if slide.shapes.title:
            title_txt = (slide.shapes.title.text or "").strip()
            if title_txt:
                out.append(_add_source_top_level({
                    "element_id": f"{file_sha1}-slide{s_idx}-title",
                    "type": "Title",
                    "text": title_txt,
                    "metadata": {
                        "file_name": path.name,
                        "file_path": str(path),
                        "file_sha1": file_sha1,
                        "source": "pptx_local",
                        "mime_type": detect_mime(path),
                        "page_number": s_idx,
                        "slide_number": s_idx,
                        "languages": "vie",
                    },
                }, file_sha1=file_sha1, file_path=path))

        # Other shapes
        for shp_i, shp in enumerate(slide.shapes, start=1):
            try:
                if getattr(shp, "has_text_frame", False) and shp.text_frame and shp.text:
                    txt = (shp.text or "").strip()
                    if txt and txt != title_txt:
                        out.append(_add_source_top_level({
                            "element_id": f"{file_sha1}-slide{s_idx}-shape{shp_i}-p",
                            "type": "NarrativeText",
                            "text": txt,
                            "metadata": {
                                "file_name": path.name,
                                "file_path": str(path),
                                "file_sha1": file_sha1,
                                "source": "pptx_local",
                                "mime_type": detect_mime(path),
                                "page_number": s_idx,
                                "slide_number": s_idx,
                                "languages": "vie",
                            },
                        }, file_sha1=file_sha1, file_path=path))
                if getattr(shp, "has_table", False) and shp.table:
                    md = _pptx_extract_markdown_table(shp.table)
                    if md.strip():
                        out.append(_add_source_top_level({
                            "element_id": f"{file_sha1}-slide{s_idx}-shape{shp_i}-tbl",
                            "type": "Table",
                            "text": "[BẢNG]\n" + md,
                            "metadata": {
                                "file_name": path.name,
                                "file_path": str(path),
                                "file_sha1": file_sha1,
                                "source": "pptx_local",
                                "mime_type": detect_mime(path),
                                "page_number": s_idx,
                                "slide_number": s_idx,
                                "languages": "vie",
                                # keep a tiny HTML-ish echo for later if needed
                                "text_as_html": None,
                            },
                        }, file_sha1=file_sha1, file_path=path))
                if getattr(shp, "shape_type", None) and str(getattr(shp, "shape_type")).lower().find("picture") >= 0:
                    # Picture placeholder (no caption extraction here)
                    out.append(_add_source_top_level({
                        "element_id": f"{file_sha1}-slide{s_idx}-shape{shp_i}-img",
                        "type": "Image",
                        "text": "[HÌNH] Ảnh/biểu đồ trên slide",
                        "metadata": {
                            "file_name": path.name,
                            "file_path": str(path),
                            "file_sha1": file_sha1,
                            "source": "pptx_local",
                            "mime_type": detect_mime(path),
                            "page_number": s_idx,
                            "slide_number": s_idx,
                            "languages": "vie",
                        },
                    }, file_sha1=file_sha1, file_path=path))
            except Exception:
                continue
    return out

# ----------------- Core -----------------
def preprocess_file(path: Path, out_jsonl: Path, max_pdf_pages_per_chunk: int) -> int:
    suffix = path.suffix.lower()
    if suffix not in SUPPORTED_EXTS:
        return 0
    file_sha1 = sha1_of_file(path)
    lines_written = 0
    file_size = path.stat().st_size
    timeout_pair = _compute_upload_timeouts(file_size)

    if suffix == ".pdf":
        for chunk_bytes, start_idx, end_idx in chunk_pdf_bytes(path, max_pages=max_pdf_pages_per_chunk):
            params = build_partition_params_for(path)
            resp_json = _post_to_unstructured_bytes(chunk_bytes, path.name, params, timeout_pair=timeout_pair)
            offset = max(start_idx - 1, 0) if end_idx != -1 else 0
            with out_jsonl.open("a", encoding="utf-8") as out:
                for el in resp_json:
                    norm = _normalize_element(el, file_sha1=file_sha1, file_path=path, abs_page_offset=offset)
                    out.write(json.dumps(norm, ensure_ascii=False) + "\n")
                    lines_written += 1
    elif suffix in {".doc", ".docx"}:
        params = build_partition_params_for(path)
        file_bytes = path.read_bytes()
        resp_json = _post_to_unstructured_bytes(file_bytes, path.name, params, timeout_pair=timeout_pair)
        with out_jsonl.open("a", encoding="utf-8") as out:
            for el in resp_json:
                norm = _normalize_element(el, file_sha1=file_sha1, file_path=path, abs_page_offset=0)
                out.write(json.dumps(norm, ensure_ascii=False) + "\n")
                lines_written += 1
    elif suffix in PPT_EXTS:
        # PPT/PPTX handled OFFLINE
        try:
            rows = _pptx_process(path, file_sha1=file_sha1)
        except Exception as e:
            print(f"⚠️  Failed PPT/PPTX local parse {path.name}: {e}", file=sys.stderr)
            return 0
        with out_jsonl.open("a", encoding="utf-8") as out:
            for norm in rows:
                out.write(json.dumps(norm, ensure_ascii=False) + "\n")
                lines_written += 1
    else:
        return 0
    return lines_written

def preprocess_folder(data_dir: str, out_jsonl: str, max_pdf_pages_per_chunk: int) -> Tuple[int, int]:
    data_path = Path(data_dir)
    out_path = Path(out_jsonl)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    existing_hashes = load_existing_hashes(out_path)

    seen_files = 0
    new_files = 0
    for p in sorted(data_path.rglob("*")):
        if not p.is_file(): continue
        if p.suffix.lower() not in SUPPORTED_EXTS: continue
        seen_files += 1
        fsha1 = sha1_of_file(p)
        if fsha1 in existing_hashes:
            continue
        try:
            written = preprocess_file(p, out_path, max_pdf_pages_per_chunk=max_pdf_pages_per_chunk)
            if written > 0:
                new_files += 1
                existing_hashes.add(fsha1)
                print(f"Appended {written} elements from: {p.name}")
        except Exception as e:
            print(f"⚠️  Failed {p.name}: {e}", file=sys.stderr)
    return seen_files, new_files

# ----------------- CLI -----------------
def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--data-dir", required=True, help="Folder with DOC/DOCX/PDF/PPT/PPTX")
    parser.add_argument("--out-jsonl", required=True, help="Output JSONL to append elements")
    parser.add_argument("--max-pdf-pages-per-chunk", type=int,
                        default=int(os.environ.get("MAX_PDF_PAGES_PER_CHUNK", DEFAULT_MAX_PDF_PAGES_PER_CHUNK)))
    args = parser.parse_args()
    n_seen, n_new = preprocess_folder(args.data_dir, args.out_jsonl, max_pdf_pages_per_chunk=args.max_pdf_pages_per_chunk)
    print(f"✅ Scanned {n_seen} files, appended {n_new} new files to {args.out_jsonl}")

if __name__ == "__main__":
    main()
